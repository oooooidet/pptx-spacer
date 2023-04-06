import pptx
from pptx.enum.shapes import PP_PLACEHOLDER
import glob
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
import re
import argparse
import os

parser = argparse.ArgumentParser()
parser.add_argument("--src-dir", type=str, default="./input", help="directory path for source pptx")
parser.add_argument("--tgt-dir", type=str, default="./output", help="directory path for target pptx")

opt = parser.parse_args()

dir_src = opt.src_dir
dir_tgt = opt.tgt_dir

pptx_files = [f for f in glob.glob(os.path.join(dir_src, "*"), recursive=True) if '.pptx' in f and '~' not in f]


def get_body_ph(placeholders):
    return list(filter(lambda ph: ph.element.ph_type == PP_PLACEHOLDER.BODY, placeholders))[0]


def refine_text(text):
    text = text.replace("\v", "")
    text = re.sub(r"([a-zA-Z0-9])([ぁ-んァ-ン一-龥])", r"\1 \2", text)
    text = re.sub(r"([ぁ-んァ-ン一-龥])([a-zA-Z0-9])", r"\1 \2", text)
    text = re.sub(r"(\w):(\w)", r"\1: \2", text)
    text = re.sub(r"([\w!?])(\()", r"\1 \2", text)
    text = re.sub(r"(\))([\w!?])", r"\1 \2", text)
    return text


def refine_txbody(txbody):
    for ap in txbody.xpath("a:p"):
        if len(ap.text) == 0:
            continue
        text = refine_text(ap.text)
        ar_list = ap.findall("a:r", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        for ar in ar_list:
            ar_words = ar.text.split(" ")
            ar_words = [refine_text(word) + " " if i < len(ar_words) - 1 else refine_text(word) for i, word in enumerate(ar_words)]
            ar.text = ""
            for ar_word in ar_words:
                if not text.startswith(ar_word):
                    if text.startswith(" "):
                        ar_word = " " + ar_word
                    elif text.endswith(" "):
                        ar_word = ar_word + " "
                text = text[len(ar_word):]
                ar.text = ar.text + ar_word


def refine_shapes(shapes):
    for shp in shapes:
        # shp == GROUP であれば、再帰呼び出し
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            refine_shapes(shp.shapes)
        # shp == TEXT であれば、普通に校正
        elif shp.has_text_frame and shp.text:
            refine_txbody(shp.element.txBody)
        # shp == TABLE であれば、セルごとに取得して校正
        elif shp.has_table:
            tbl = shp.table
            for row in tbl.rows:
                for cell in row.cells:
                    refine_txbody(cell.text_frame._txBody)


for f in pptx_files:
    p = pptx.Presentation(f)
    print(os.path.basename(f))
    for i, slide in enumerate(p.slides, start=1):
        print("slide #" + str(i))
        # Slide
        refine_shapes(slide.shapes)
        # Note
        note = slide.notes_slide
        refine_txbody(get_body_ph(note.placeholders).element.txBody)
    p.save(os.path.join(dir_tgt, os.path.basename(f)))
